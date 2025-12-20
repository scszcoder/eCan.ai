"""
File Extractors for LightRAG Knowledge Base

This module provides file extraction utilities specifically for LightRAG,
extracting content in formats optimized for intelligent chunking.

Supported formats:
- Excel (.xlsx, .xls) -> Markdown tables
- CSV/TSV -> Preserved as-is (handled by chunker)
- Other formats can be added as needed
"""

import io
import os
from typing import Optional
from utils.logger_helper import logger_helper as logger


def extract_excel_as_markdown(file_bytes: bytes, filename: str) -> str:
    """
    Extract Excel file content with optional Pandas preprocessing.
    
    If ENABLE_PANDAS_PREPROCESSING=1, uses Pandas for semantic enhancement.
    Otherwise, uses basic tab-separated format.
    
    Format (basic):
    Sheet: SheetName
    Col1    Col2    Col3
    Val1    Val2    Val3
    
    Format (Pandas enhanced):
    # Table: filename - SheetName
    **Metadata**: ...
    **Statistics**: ...
    **Data**: Markdown table
    """
    # Check if Pandas preprocessing is enabled
    use_pandas = os.getenv('ENABLE_PANDAS_PREPROCESSING', '0') == '1'
    
    if use_pandas:
        try:
            from knowledge.table_preprocessor import preprocess_excel_with_pandas
            result = preprocess_excel_with_pandas(file_bytes, filename)
            if result:
                logger.info(f"[FileExtractor] ✅ Used Pandas preprocessing for {filename}")
                return result
            else:
                logger.warning(f"[FileExtractor] Pandas preprocessing failed, falling back to basic extraction")
        except Exception as e:
            logger.warning(f"[FileExtractor] Pandas preprocessing error: {e}, falling back to basic extraction")
    
    # Basic extraction (original logic)
    try:
        import openpyxl
        logger.info(f"[FileExtractor] Extracting Excel file: {filename}")
        
        # Load workbook in read-only mode for better performance
        wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
        logger.debug(f"[FileExtractor] Loaded workbook with {len(wb.worksheets)} sheets")
        
        # Extract each sheet
        sheet_blocks = []
        
        for sheet in wb.worksheets:
            # Sheet title marker (parsed by chunker)
            sheet_block = [f"Sheet: {sheet.title}"]
            
            # Get all rows
            rows = list(sheet.iter_rows(values_only=True))
            if not rows:
                logger.debug(f"[FileExtractor] Sheet '{sheet.title}' is empty, skipping")
                continue
            
            logger.debug(f"[FileExtractor] Processing sheet '{sheet.title}' with {len(rows)} rows")
            
            # Find the maximum column index with non-empty data across ALL rows
            # This fixes the "16384 empty columns" issue
            max_col_idx = 0
            for row in rows:
                for i in range(len(row) - 1, -1, -1):
                    if row[i] is not None and str(row[i]).strip():
                        max_col_idx = max(max_col_idx, i)
                        break
            
            # If no data found, skip this sheet
            if max_col_idx == 0 and not any(row[0] for row in rows if row):
                logger.debug(f"[FileExtractor] Sheet '{sheet.title}' has no data, skipping")
                continue
            
            logger.debug(f"[FileExtractor] Sheet '{sheet.title}' max column index: {max_col_idx} (trimmed from {len(rows[0]) if rows else 0} columns)")
            
            # Convert rows to tab-separated strings, trimming to max_col_idx + 1
            for row in rows:
                # Skip completely empty rows
                if not any(cell is not None and str(cell).strip() for cell in row[:max_col_idx + 1]):
                    continue
                
                # Trim row to max_col_idx + 1 and convert to string
                trimmed_row = row[:max_col_idx + 1]
                row_str = "\t".join([str(cell) if cell is not None else "" for cell in trimmed_row])
                sheet_block.append(row_str)
            
            if len(sheet_block) > 1: # Only add if has data
                sheet_blocks.append("\n".join(sheet_block))
        
        # Separate sheets with double newlines
        result = "\n\n".join(sheet_blocks)
        logger.info(f"[FileExtractor] ✅ Successfully extracted {len(sheet_blocks)} sheets")
        return result
        
    except ImportError as e:
        logger.error(f"[FileExtractor] ❌ openpyxl not installed: {e}")
        return f"[Excel extraction error: openpyxl not installed]"
    except Exception as e:
        logger.error(f"[FileExtractor] ❌ Failed to extract Excel file '{filename}': {e}")
        return f"[Excel extraction error: {e}]"


def extract_docx_with_tables(file_bytes: bytes, filename: str) -> str:
    """
    Extract DOCX content including tables as Markdown.
    
    Extracts both paragraphs and tables, preserving document structure.
    Tables are converted to Markdown format for intelligent chunking.
    
    Args:
        file_bytes: Raw file bytes
        filename: Original filename
        
    Returns:
        Mixed content with text paragraphs and Markdown tables
    """
    try:
        import docx
        from docx.table import Table
        from docx.text.paragraph import Paragraph
        
        logger.info(f"[FileExtractor] Extracting DOCX file: {filename}")
        
        doc = docx.Document(io.BytesIO(file_bytes))
        content_parts = []
        table_index = 0
        
        # Iterate through document elements in order to preserve structure
        for element in doc.element.body:
            # Paragraph
            if element.tag.endswith('p'):
                para = Paragraph(element, doc)
                if para.text.strip():
                    content_parts.append(para.text.strip())
            
            # Table
            elif element.tag.endswith('tbl'):
                table = Table(element, doc)
                table_index += 1
                
                # Convert table to Markdown
                markdown_table = _convert_docx_table_to_markdown(
                    table, 
                    f"Table {table_index}"
                )
                if markdown_table:
                    content_parts.append(markdown_table)
                    logger.debug(f"[FileExtractor] Extracted table {table_index} ({len(table.rows)} rows, {len(table.columns)} columns)")
        
        result = "\n\n".join(content_parts)
        logger.info(f"[FileExtractor] ✅ Extracted DOCX with {table_index} tables, {len(content_parts)} content blocks")
        return result
        
    except ImportError:
        logger.error("[FileExtractor] ❌ python-docx not installed")
        return "[DOCX extraction error: python-docx not installed]"
    except Exception as e:
        logger.error(f"[FileExtractor] ❌ Failed to extract DOCX '{filename}': {e}")
        return f"[DOCX extraction error: {e}]"


def _convert_docx_table_to_markdown(table, title: str = "Table") -> str:
    """
    Convert DOCX table to Markdown format.
    
    Args:
        table: python-docx Table object
        title: Table title for heading
        
    Returns:
        Markdown formatted table string
    """
    try:
        if not table.rows:
            return ""
        
        lines = [f"## {title}", ""]
        
        # First row as header
        header_cells = [cell.text.strip() for cell in table.rows[0].cells]
        if not any(header_cells):  # Empty header row
            return ""
        
        lines.append("| " + " | ".join(header_cells) + " |")
        lines.append("|" + "|".join(["---"] * len(header_cells)) + "|")
        
        # Data rows
        for row in table.rows[1:]:
            cells = [cell.text.strip() for cell in row.cells]
            # Pad to match header length
            while len(cells) < len(header_cells):
                cells.append("")
            lines.append("| " + " | ".join(cells[:len(header_cells)]) + " |")
        
        return "\n".join(lines)
    except Exception as e:
        logger.warning(f"[FileExtractor] Failed to convert DOCX table to Markdown: {e}")
        return ""


def extract_pptx_with_tables(file_bytes: bytes, filename: str) -> str:
    """
    Extract PPTX content including tables as Markdown.
    
    Extracts text and tables from each slide, preserving slide structure.
    Tables are converted to Markdown format for intelligent chunking.
    
    Args:
        file_bytes: Raw file bytes
        filename: Original filename
        
    Returns:
        Slide content with text and Markdown tables
    """
    try:
        from pptx import Presentation
        
        logger.info(f"[FileExtractor] Extracting PPTX file: {filename}")
        
        prs = Presentation(io.BytesIO(file_bytes))
        slide_contents = []
        total_tables = 0
        
        for slide_idx, slide in enumerate(prs.slides, 1):
            slide_parts = [f"# Slide {slide_idx}"]
            table_count = 0
            
            for shape in slide.shapes:
                # Text content
                if hasattr(shape, "text") and shape.text.strip():
                    slide_parts.append(shape.text.strip())
                
                # Table
                if shape.has_table:
                    table_count += 1
                    total_tables += 1
                    markdown_table = _convert_pptx_table_to_markdown(
                        shape.table,
                        f"Table {table_count}"
                    )
                    if markdown_table:
                        slide_parts.append(markdown_table)
                        logger.debug(f"[FileExtractor] Extracted table {table_count} from slide {slide_idx} ({len(shape.table.rows)} rows)")
            
            if len(slide_parts) > 1:  # Has content beyond slide number
                slide_contents.append("\n\n".join(slide_parts))
        
        result = "\n\n---\n\n".join(slide_contents)
        logger.info(f"[FileExtractor] ✅ Extracted PPTX with {total_tables} tables from {len(prs.slides)} slides")
        return result
        
    except ImportError:
        logger.error("[FileExtractor] ❌ python-pptx not installed")
        return "[PPTX extraction error: python-pptx not installed]"
    except Exception as e:
        logger.error(f"[FileExtractor] ❌ Failed to extract PPTX '{filename}': {e}")
        return f"[PPTX extraction error: {e}]"


def _convert_pptx_table_to_markdown(table, title: str = "Table") -> str:
    """
    Convert PPTX table to Markdown format.
    
    Args:
        table: python-pptx Table object
        title: Table title for heading
        
    Returns:
        Markdown formatted table string
    """
    try:
        if not table.rows:
            return ""
        
        lines = [f"## {title}", ""]
        
        # First row as header
        header_cells = [cell.text.strip() for cell in table.rows[0].cells]
        if not any(header_cells):  # Empty header row
            return ""
        
        lines.append("| " + " | ".join(header_cells) + " |")
        lines.append("|" + "|".join(["---"] * len(header_cells)) + "|")
        
        # Data rows
        for row in table.rows[1:]:
            cells = [cell.text.strip() for cell in row.cells]
            # Pad to match header length
            while len(cells) < len(header_cells):
                cells.append("")
            lines.append("| " + " | ".join(cells[:len(header_cells)]) + " |")
        
        return "\n".join(lines)
    except Exception as e:
        logger.warning(f"[FileExtractor] Failed to convert PPTX table to Markdown: {e}")
        return ""


def extract_csv_content(file_bytes: bytes, filename: str) -> str:
    """
    Extract CSV/TSV content.
    
    CSV files are returned as-is since the advanced chunker can handle
    CSV format directly and convert it to Markdown tables.
    
    Args:
        file_bytes: Raw file bytes
        filename: Original filename
        
    Returns:
        CSV content as string
    """
    try:
        logger.info(f"[FileExtractor] Extracting CSV file: {filename}")
        result = file_bytes.decode("utf-8", errors="ignore")
        logger.info(f"[FileExtractor] ✅ Successfully extracted CSV content ({len(result)} chars)")
        return result
    except Exception as e:
        logger.error(f"[FileExtractor] ❌ Failed to extract CSV file '{filename}': {e}")
        return f"[CSV extraction error: {e}]"


def extract_file_for_lightrag(file_bytes: bytes, filename: str) -> str:
    """
    Main entry point for file extraction optimized for LightRAG.
    
    This function routes to the appropriate extractor based on file extension.
    
    Args:
        file_bytes: Raw file bytes
        filename: Original filename
        
    Returns:
        Extracted content in format optimized for LightRAG chunking
    """
    import os
    
    ext = os.path.splitext(filename)[1].lower()
    logger.info(f"[FileExtractor] Processing file: {filename} (type: {ext})")
    
    # Excel files -> Markdown tables
    if ext in (".xlsx", ".xls"):
        return extract_excel_as_markdown(file_bytes, filename)
    
    # Word documents -> Text + Markdown tables
    elif ext in (".docx",):
        return extract_docx_with_tables(file_bytes, filename)
    
    # PowerPoint -> Text + Markdown tables
    elif ext in (".pptx", ".ppt"):
        return extract_pptx_with_tables(file_bytes, filename)
    
    # CSV/TSV -> As-is (chunker handles it)
    elif ext in (".csv", ".tsv"):
        return extract_csv_content(file_bytes, filename)
    
    # Other formats not supported by this module
    else:
        logger.warning(f"[FileExtractor] ⚠️  Unsupported format: {ext} for file '{filename}'")
        return f"[Unsupported format for LightRAG extraction: {ext}]"


# Backward compatibility function
def extract_excel_to_markdown(file_bytes: bytes) -> str:
    """
    Legacy function name for backward compatibility.
    
    Args:
        file_bytes: Raw Excel file bytes
        
    Returns:
        Markdown-formatted table content
    """
    return extract_excel_as_markdown(file_bytes, "unknown.xlsx")
