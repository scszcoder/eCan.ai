import io
import os

def extract_file_text(file_bytes: bytes, filename: str) -> str:
    ext = os.path.splitext(filename)[1].lower()
    if ext == ".pdf":
        try:
            import PyPDF2
            with io.BytesIO(file_bytes) as buf:
                reader = PyPDF2.PdfReader(buf)
                return "\n".join(page.extract_text() for page in reader.pages if page.extract_text())
        except Exception as e:
            return f"[PDF extraction error: {e}]"
    elif ext in (".xlsx", ".xls"):
        try:
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True)
            text = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    text.append("\t".join([str(cell) if cell is not None else "" for cell in row]))
            return "\n".join(text)
        except Exception as e:
            return f"[Excel extraction error: {e}]"
    elif ext in (".pptx", ".ppt"):
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_bytes))
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)
        except Exception as e:
            return f"[PPT extraction error: {e}]"
    elif ext in (".docx",):
        try:
            import docx
            doc = docx.Document(io.BytesIO(file_bytes))
            return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
        except Exception as e:
            return f"[Word extraction error: {e}]"
    elif ext in (".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff"):
        try:
            from PIL import Image
            import pytesseract
            img = Image.open(io.BytesIO(file_bytes))
            text = pytesseract.image_to_string(img)
            return text
        except Exception as e:
            return f"[Image OCR error: {e}]"
    elif ext in (".txt", ".csv", ".md"):
        try:
            return file_bytes.decode("utf-8", errors="ignore")
        except Exception as e:
            return f"[Text extraction error: {e}]"
    elif ext in (".mp4", ".avi", ".mov", ".wmv", ".flv", ".webm", ".mkv"):
        try:
            import cv2
            from PIL import Image
            import pytesseract

            # Extract text from a few frames only (e.g. 10 evenly spread)
            with io.BytesIO(file_bytes) as buf:
                temp_file = "temp_video_file"
                with open(temp_file, "wb") as f:
                    f.write(buf.read())
                cap = cv2.VideoCapture(temp_file)
                total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
                texts = []
                frame_indexes = list(range(0, total_frames, max(1, total_frames // 10)))
                for idx in frame_indexes:
                    cap.set(cv2.CAP_PROP_POS_FRAMES, idx)
                    ret, frame = cap.read()
                    if ret:
                        # Convert to PIL for pytesseract
                        im_pil = Image.fromarray(cv2.cvtColor(frame, cv2.COLOR_BGR2RGB))
                        texts.append(pytesseract.image_to_string(im_pil))
                cap.release()
                os.remove(temp_file)
            return "\n".join(texts)
        except Exception as e:
            return f"[Video extraction error: {e}]"
    else:
        return "[Unsupported file type for text extraction]"
